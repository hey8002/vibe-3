from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("project_structure_diagram.pptx")
FONT = "Malgun Gothic"

COLORS = {
    "navy": RGBColor(18, 34, 51),
    "blue": RGBColor(37, 99, 235),
    "teal": RGBColor(13, 148, 136),
    "green": RGBColor(22, 163, 74),
    "orange": RGBColor(234, 88, 12),
    "yellow": RGBColor(245, 158, 11),
    "purple": RGBColor(124, 58, 237),
    "red": RGBColor(220, 38, 38),
    "gray": RGBColor(100, 116, 139),
    "ink": RGBColor(15, 23, 42),
    "soft": RGBColor(248, 250, 252),
    "white": RGBColor(255, 255, 255),
}


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(226, 232, 240)


def write_text(shape, text: str, size: int = 14, bold: bool = False, color: RGBColor | None = None, align=PP_ALIGN.CENTER) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]


def add_header(slide, title: str, subtitle: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.72))
    set_fill(bar, COLORS["navy"])
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.13), Inches(8.5), Inches(0.45))
    write_text(title_box, title, 22, True, COLORS["white"], PP_ALIGN.LEFT)

    sub_box = slide.shapes.add_textbox(Inches(9.0), Inches(0.18), Inches(3.8), Inches(0.35))
    write_text(sub_box, subtitle, 10, False, RGBColor(203, 213, 225), PP_ALIGN.RIGHT)


def add_footer(slide, page: int) -> None:
    footer = slide.shapes.add_textbox(Inches(10.5), Inches(7.12), Inches(2.35), Inches(0.25))
    write_text(footer, f"day3_rpa 프로젝트 구조 | {page}", 9, False, COLORS["gray"], PP_ALIGN.RIGHT)


def add_box(slide, x: float, y: float, w: float, h: float, text: str, fill: RGBColor, size: int = 12, bold: bool = True, font_color: RGBColor | None = None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(box, fill)
    write_text(box, text, size, bold, font_color or COLORS["ink"])
    return box


def add_textbox(slide, x: float, y: float, w: float, h: float, text: str, size: int = 13, bold: bool = False, color: RGBColor | None = None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    for index, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.bold = bold and index == 0
        p.font.color.rgb = color or COLORS["ink"]
        p.space_after = Pt(2)
    return box


def add_bullets(slide, x: float, y: float, w: float, h: float, title: str, bullets: list[str], color: RGBColor):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(box, COLORS["white"])
    box.line.color.rgb = color
    box.line.width = Pt(1.3)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.13)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = color
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = FONT
        p.font.size = Pt(10.5)
        p.font.color.rgb = COLORS["ink"]
        p.space_after = Pt(1.5)
    return box


def connect(slide, left, right, color: RGBColor | None = None) -> None:
    start_x = left.left + left.width
    start_y = left.top + left.height // 2
    end_x = right.left
    end_y = right.top + right.height // 2
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    line.line.color.rgb = color or COLORS["gray"]
    line.line.width = Pt(1.5)


# 1. Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
set_fill(bg, COLORS["navy"])
bg.line.fill.background()
add_textbox(slide, 0.75, 1.05, 9.5, 0.75, "day3_rpa 프로젝트 구조 상세 도식화", 34, True, COLORS["white"])
add_textbox(slide, 0.78, 1.95, 10.5, 0.55, "공공직군 행정업무 슈퍼앱: Frontend + Backend + SQLite + RAG/Excel/News", 16, False, RGBColor(203, 213, 225))
add_box(slide, 0.8, 3.0, 2.5, 0.7, "React + TypeScript + Vite", COLORS["blue"], 13, True, COLORS["white"])
add_box(slide, 3.55, 3.0, 2.4, 0.7, "Python + FastAPI + uv", COLORS["teal"], 13, True, COLORS["white"])
add_box(slide, 6.2, 3.0, 1.8, 0.7, "SQLite", COLORS["green"], 13, True, COLORS["white"])
add_box(slide, 8.25, 3.0, 2.25, 0.7, "OpenAI API + RAG", COLORS["orange"], 13, True, COLORS["white"])
add_textbox(slide, 0.8, 5.75, 8.0, 0.45, "생성 파일: docs/project_structure_diagram.pptx", 13, False, RGBColor(203, 213, 225))
add_footer(slide, 1)

# 2. Overall Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "1. 전체 아키텍처", "Client - API - Service - Data")
user = add_box(slide, 0.55, 2.8, 1.65, 0.75, "사용자\n브라우저", COLORS["yellow"], 14)
fe = add_box(slide, 2.9, 2.35, 2.35, 1.65, "Frontend\nReact / TS / Vite\nGitHub Pages 대상", COLORS["blue"], 13, True, COLORS["white"])
api = add_box(slide, 5.95, 2.35, 2.35, 1.65, "Backend API\nFastAPI / uvicorn\nCORS + .env 로딩", COLORS["teal"], 13, True, COLORS["white"])
excel = add_box(slide, 9.0, 1.15, 2.0, 0.75, "Excel 자동화", RGBColor(219, 234, 254), 12)
rag = add_box(slide, 9.0, 2.2, 2.0, 0.75, "RAG 챗봇", RGBColor(255, 237, 213), 12)
news = add_box(slide, 9.0, 3.25, 2.0, 0.75, "뉴스 수집", RGBColor(220, 252, 231), 12)
db = add_box(slide, 11.55, 2.1, 1.3, 0.85, "SQLite\napp.db", COLORS["green"], 12, True, COLORS["white"])
files = add_box(slide, 11.55, 3.35, 1.3, 0.95, "업로드/결과\n파일", COLORS["gray"], 11, True, COLORS["white"])
for a, b in [(user, fe), (fe, api), (api, excel), (api, rag), (api, news), (excel, db), (rag, db), (news, db), (excel, files), (rag, files)]:
    connect(slide, a, b)
add_textbox(slide, 0.75, 5.55, 11.8, 0.55, "핵심 흐름: 브라우저 UI → /api REST 호출 → FastAPI 라우터 → 서비스 계층 처리 → SQLite/파일 저장소 반영", 15, True)
add_footer(slide, 2)

# 3. Directory Structure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "2. 루트 디렉터리 구조", "실제 파일 기준")
tree = """day3_rpa/
├─ frontend/            React + TypeScript + Vite 클라이언트
│  ├─ src/app/           App shell, 메뉴, API URL 설정
│  ├─ src/pages/         Home, Schedule, Excel, Chatbot, News 화면
│  └─ src/shared/api/    fetch 기반 API 클라이언트
├─ backend/             FastAPI 서버
│  ├─ app/api/routes/    REST 엔드포인트
│  ├─ app/services/      Excel, News, Chatbot 비즈니스 로직
│  ├─ app/schemas/       요청/응답 Pydantic 스키마
│  ├─ app/core/          DB 경로, 설정
│  └─ data/              SQLite DB, 업로드/결과 파일 저장소
├─ docs/                PRD, Architecture, Operation, 발표자료
├─ .github/workflows/   GitHub Pages 배포 자동화
└─ .codex/.agent/       Codex 스킬 및 작업 설정"""
tree_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.12), Inches(6.55), Inches(5.75))
set_fill(tree_box, COLORS["soft"])
write_text(tree_box, tree, 12, False, COLORS["ink"], PP_ALIGN.LEFT)
for paragraph in tree_box.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.name = "Consolas"
add_bullets(slide, 7.55, 1.25, 4.95, 1.3, "분리 기준", ["frontend와 backend를 명확히 분리", "docs는 요구사항/운영/아키텍처 문서 보관", "data와 .env는 Git 제외 대상"], COLORS["blue"])
add_bullets(slide, 7.55, 2.85, 4.95, 1.3, "런타임 산출물", ["backend/data/app.db", "backend/data/uploads, outputs", "backend/data/chatbot_uploads"], COLORS["green"])
add_bullets(slide, 7.55, 4.45, 4.95, 1.45, "주의 지점", ["실제 API 키는 backend/.env에만 보관", ".env.example만 커밋", "업로드 파일과 DB는 배포 산출물이 아님"], COLORS["orange"])
add_footer(slide, 3)

# 4. Frontend
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "3. Frontend 구조", "React 화면과 API 클라이언트")
main = add_box(slide, 0.7, 1.25, 2.2, 0.75, "src/main.tsx\nReact 진입점", COLORS["blue"], 12, True, COLORS["white"])
app = add_box(slide, 3.55, 1.25, 2.35, 0.75, "src/app/App.tsx\n메뉴/페이지 전환/API URL", COLORS["blue"], 12, True, COLORS["white"])
pages = [
    ("HomePage", "상태 확인\nAPI URL 설정"),
    ("SchedulePage", "일정 화면"),
    ("ExcelPage", "엑셀 업로드/분리/병합"),
    ("ChatbotPage", "RAG 업로드/학습/질문"),
    ("NewsPage", "정책뉴스 수집/목록"),
]
page_boxes = []
for index, (name, desc) in enumerate(pages):
    page_boxes.append(add_box(slide, 6.55, 0.8 + index * 1.05, 2.15, 0.78, f"{name}\n{desc}", RGBColor(219, 234, 254), 10))
api = add_box(slide, 10.05, 1.25, 2.25, 4.0, "src/shared/api/\nbase.ts\nhealth.ts\nexcel.ts\nchatbot.ts\nnews.ts", COLORS["soft"], 12)
connect(slide, main, app)
for page in page_boxes:
    connect(slide, app, page, COLORS["blue"])
    connect(slide, page, api, COLORS["blue"])
add_textbox(slide, 0.75, 6.25, 11.8, 0.45, "프론트는 App 상태(activePage)로 화면을 전환하며, API base URL은 localStorage 또는 VITE_API_BASE_URL로 설정됩니다.", 13)
add_footer(slide, 4)

# 5. Backend
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "4. Backend 구조", "FastAPI 라우터 - 서비스 - 저장소")
main = add_box(slide, 0.65, 1.05, 2.25, 0.85, "app/main.py\nFastAPI 앱 / CORS / .env", COLORS["teal"], 11, True, COLORS["white"])
routes = add_box(slide, 3.55, 0.9, 2.1, 1.25, "api/routes\nhealth.py\nexcel.py\nchatbot.py\nnews.py\nschedules.py", RGBColor(204, 251, 241), 10)
schemas = add_box(slide, 3.55, 2.65, 2.1, 0.95, "schemas\n요청/응답 모델", RGBColor(224, 242, 254), 12)
services = add_box(slide, 6.4, 0.9, 2.4, 2.7, "services\nexcel_processor/store/workbook\nchatbot_pipeline/processor/store\nnews_collector/scheduler/store", RGBColor(255, 237, 213), 10)
core = add_box(slide, 9.45, 0.9, 1.95, 0.95, "core\nconfig.py\ndatabase.py", COLORS["soft"], 11)
data = add_box(slide, 9.45, 2.35, 1.95, 1.25, "data\napp.db\nuploads\noutputs\nchatbot_uploads", COLORS["green"], 10, True, COLORS["white"])
external = add_box(slide, 11.85, 1.8, 1.15, 1.0, "외부\nOpenAI\n정책브리핑", COLORS["orange"], 9, True, COLORS["white"])
for a, b in [(main, routes), (routes, schemas), (routes, services), (services, core), (services, data), (services, external)]:
    connect(slide, a, b, COLORS["teal"])
add_bullets(slide, 0.7, 4.55, 3.75, 1.35, "실행", ["cd backend", "uv run uvicorn app.main:app --reload", "OPENAI_API_KEY는 backend/.env에서 로드"], COLORS["teal"])
add_bullets(slide, 4.8, 4.55, 3.75, 1.35, "주요 의존성", ["FastAPI, uvicorn, python-dotenv", "openpyxl, pandas", "LangChain, OpenAI, pypdf", "BeautifulSoup4, APScheduler"], COLORS["orange"])
add_bullets(slide, 8.9, 4.55, 3.65, 1.35, "저장소", ["SQLite 단일 DB", "업로드 원본 및 결과 파일 분리", "Git에는 산출물 제외"], COLORS["green"])
add_footer(slide, 5)

# 6. API Map
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "5. API 엔드포인트 맵", "prefix: /api")
columns = [
    ("Health", ["GET /health", "GET /db/health"], COLORS["gray"]),
    ("Excel", ["GET /excel/jobs", "GET /excel/stored", "GET /excel/stored/{id}", "POST /excel/upload", "POST /excel/split", "POST /excel/merge", "GET /excel/download/{file}", "DELETE /excel/stored/{id}", "DELETE /excel/stored"], COLORS["blue"]),
    ("Chatbot", ["GET /chatbot/documents", "POST /chatbot/documents/upload", "POST /chatbot/train", "POST /chatbot/ask"], COLORS["orange"]),
    ("News", ["GET /news", "GET /news/jobs", "POST /news/collect"], COLORS["green"]),
    ("Schedules", ["GET /schedules"], COLORS["purple"]),
]
for index, (title, items, color) in enumerate(columns):
    add_bullets(slide, 0.45 + index * 2.55, 1.15, 2.25, 5.25 if title == "Excel" else 4.2, title, items, color)
add_textbox(slide, 0.6, 6.65, 12, 0.35, "프론트 API 클라이언트는 src/shared/api/*에 있고, base.ts가 실제 API Base URL을 조합합니다.", 12)
add_footer(slide, 6)

# 7. Feature Flows
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "6. 기능별 데이터 흐름", "Excel / RAG / News")
flows = [
    ("Excel 자동화", COLORS["blue"], ["파일 업로드", "컬럼 선택", "분리/병합 처리", "DB 저장/결과 파일", "표 조회/다운로드"]),
    ("RAG 챗봇", COLORS["orange"], ["문서 업로드", "텍스트 추출", "전처리/청크", "OpenAI 임베딩", "질문 답변/근거 표시"]),
    ("뉴스 수집", COLORS["green"], ["날짜 선택/스케줄", "BS4 크롤링", "중복 제거", "SQLite 저장", "목록 조회"]),
]
for row, (label, color, steps) in enumerate(flows):
    y = 1.35 + row * 1.75
    previous = add_box(slide, 0.55, y, 1.55, 0.65, label, color, 11, True, COLORS["white"])
    for index, step in enumerate(steps):
        current = add_box(slide, 2.55 + index * 1.95, y, 1.55, 0.65, step, COLORS["soft"], 9)
        connect(slide, previous, current, color)
        previous = current
add_footer(slide, 7)

# 8. Deploy / Config
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "7. 설정·배포·보안 구조", "환경변수와 GitHub Pages")
add_bullets(slide, 0.75, 1.1, 3.8, 1.65, "환경 파일", ["backend/.env.example: OpenAI/CORS 예시", "backend/.env: 실제 키 보관, Git 제외", "frontend/.env.example: VITE_API_BASE_URL 예시"], COLORS["orange"])
add_bullets(slide, 4.9, 1.1, 3.8, 1.65, "GitHub Pages", ["frontend 빌드 결과를 Pages에 배포", ".github/workflows/pages.yml 사용", "백엔드는 별도 서버 또는 터널 필요"], COLORS["blue"])
add_bullets(slide, 9.05, 1.1, 3.45, 1.65, "로컬 실행", ["Frontend: npm run dev", "Backend: uv run uvicorn app.main:app --reload", "DB: backend/data/app.db"], COLORS["green"])
browser = add_box(slide, 1.0, 4.0, 2.3, 0.75, "브라우저", COLORS["yellow"], 12)
pages = add_box(slide, 4.0, 4.0, 2.5, 0.75, "GitHub Pages\n정적 FE", COLORS["blue"], 12, True, COLORS["white"])
backend = add_box(slide, 7.25, 4.0, 2.5, 0.75, "Backend URL\nCloudflared 등", COLORS["teal"], 12, True, COLORS["white"])
sqlite = add_box(slide, 10.45, 4.0, 1.6, 0.75, "SQLite\n파일 저장", COLORS["green"], 12, True, COLORS["white"])
for a, b in [(browser, pages), (pages, backend), (backend, sqlite)]:
    connect(slide, a, b)
add_textbox(slide, 0.8, 6.1, 11.6, 0.5, "보안 기준: 실제 API 키와 업로드/DB 산출물은 커밋하지 않고, 예시 파일만 저장소에 포함합니다.", 13, True, COLORS["red"])
add_footer(slide, 8)

# 9. Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "8. 구조 요약 및 점검 포인트", "개발/운영 관점")
add_bullets(slide, 0.75, 1.05, 3.75, 2.0, "현재 구조의 장점", ["FE/BE 책임 분리", "API 클라이언트가 shared/api로 모여 있음", "서비스 계층이 기능별로 분리됨", "SQLite 기반 MVP 구현에 적합"], COLORS["green"])
add_bullets(slide, 4.85, 1.05, 3.75, 2.0, "주의할 점", ["docs 일부 문서 인코딩 깨짐 확인 필요", "backend/data는 백업/보호 정책 필요", "OpenAI API 키 유출 방지", "GitHub Pages는 백엔드 실행을 포함하지 않음"], COLORS["orange"])
add_bullets(slide, 8.95, 1.05, 3.6, 2.0, "다음 개선 후보", ["API 테스트 코드 추가", "DB migration 도입", "업로드 파일 크기 제한", "RAG 학습 상태 비동기 작업화"], COLORS["blue"])
add_textbox(slide, 0.9, 4.25, 11.5, 1.3, "결론: 이 프로젝트는 행정업무용 단일 웹앱입니다. 프론트는 UI와 API 연결을 담당하고, 백엔드는 Excel 처리, 뉴스 수집, RAG 챗봇, SQLite 저장을 담당합니다. 배포 시 프론트는 GitHub Pages, 백엔드는 별도 실행 환경으로 분리하는 것이 핵심입니다.", 16, True)
add_footer(slide, 9)

prs.save(OUT)
print(OUT.resolve())
